"""Test the slide engine."""
import os
from techfig.engines.slides import create_presentation


def test_create_basic_presentation(tmp_path):
    out = str(tmp_path / "basic.pptx")
    slides = [
        {"title": "Welcome", "content": "First slide\nBullet two"},
        {"title": "Methodology", "content": "- Step 1\n- Step 2\n- Step 3"},
    ]
    result = create_presentation(slides, out)
    assert result == out
    assert os.path.exists(out)
    assert os.path.getsize(out) > 1000


def test_slides_with_speaker_notes(tmp_path):
    out = str(tmp_path / "notes.pptx")
    slides = [
        {"title": "Title Slide", "content": "Hello", "notes": "Remember to greet everyone"},
        {"title": "Data", "content": "Results here", "notes": "Explain the p-value"},
    ]
    result = create_presentation(slides, out)
    assert os.path.exists(result)

    # Verify notes are in the file by reading with python-pptx
    from pptx import Presentation
    prs = Presentation(result)
    for idx, slide in enumerate(prs.slides):
        notes_slide = slide.notes_slide
        assert notes_slide.notes_text_frame.text == slides[idx]["notes"]


def test_slides_with_generated_image(tmp_path):
    """Test embedding a PNG image into a slide."""
    import pandas as pd
    from techfig.engines.figures import create_chart

    # Generate a chart to embed
    df = pd.DataFrame({"x": [1, 2, 3], "y": [4, 5, 6]})
    img_path = str(tmp_path / "chart.png")
    create_chart(data=df, chart_type="line", output_path=img_path, x_col="x", y_col="y")

    out = str(tmp_path / "img_slides.pptx")
    slides = [
        {"title": "Chart Slide", "image": img_path},
    ]
    result = create_presentation(slides, out)
    assert os.path.exists(result)
    assert os.path.getsize(result) > 5000  # Should be larger with embedded image


def test_bullet_parsing(tmp_path):
    """Bullets should become separate paragraphs."""
    out = str(tmp_path / "bullets.pptx")
    slides = [
        {"title": "Bullets", "content": "- Alpha\n- Beta\n- Gamma"},
    ]
    result = create_presentation(slides, out)

    from pptx import Presentation
    prs = Presentation(result)
    slide = prs.slides[0]
    body = slide.shapes.placeholders[1].text_frame
    # Three paragraphs: first line reuses existing paragraph, 2 added
    texts = [p.text for p in body.paragraphs if p.text]
    assert texts == ["Alpha", "Beta", "Gamma"]


def test_missing_image_warns(tmp_path):
    """Non-existent image path should create a title-only slide, not crash."""
    out = str(tmp_path / "missing_img.pptx")
    slides = [
        {"title": "Missing", "image": "/nonexistent/path.png"},
    ]
    result = create_presentation(slides, out)
    assert os.path.exists(result)
